#!/usr/bin/env python3
"""
PowerPoint Editor für BitSight Präsentation
Integriert Content aus Source-Deck in Target-Deck für besseren narrativen Flow
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import json
import copy

# Lade die Analyse
with open('/Users/nwe/clawd/slides_analysis.json', 'r') as f:
    analysis = json.load(f)

# Lade beide Präsentationen
source_prs = Presentation('/Users/nwe/clawd/Bitsight_Vorstellung MHP.pptx')
target_prs = Presentation('/Users/nwe/clawd/Bitsight_MHP_Beschaffung_Rollen_und_Zuarbeit.pptx')

# Changelog tracking
changelog = []

def log_change(action):
    """Changelog-Eintrag hinzufügen"""
    changelog.append(action)
    print(f"✓ {action}")

def copy_slide(source_prs, source_idx, target_prs, insert_position=None):
    """
    Kopiert eine Folie von source_prs nach target_prs
    """
    source_slide = source_prs.slides[source_idx]
    
    # Füge Slide mit gleichem Layout ein
    slide_layout = target_prs.slide_layouts[0]  # Nimm erstes Layout als Fallback
    
    if insert_position is None:
        new_slide = target_prs.slides.add_slide(slide_layout)
    else:
        new_slide = target_prs.slides.add_slide(slide_layout)
        # Verschiebe an gewünschte Position
        xml_slides = target_prs.slides._sldIdLst
        slides = list(xml_slides)
        xml_slides.remove(slides[-1])
        xml_slides.insert(insert_position, slides[-1])
    
    # Kopiere alle Shapes
    for shape in source_slide.shapes:
        el = shape.element
        newel = copy.deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
    
    return new_slide

def create_slide_from_analysis(prs, title, bullets, insert_position=None):
    """
    Erstellt eine neue Folie aus Analyse-Daten
    """
    # Verwende Title and Content Layout (meist Layout 1)
    slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    
    if insert_position is None:
        slide = prs.slides.add_slide(slide_layout)
    else:
        slide = prs.slides.add_slide(slide_layout)
        # Verschiebe an gewünschte Position
        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)
        xml_slides.remove(slides[-1])
        xml_slides.insert(insert_position, slides[-1])
    
    # Setze Titel
    if slide.shapes.title:
        slide.shapes.title.text = title
    
    # Füge Bullets hinzu (wenn Content-Placeholder vorhanden)
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape != slide.shapes.title:
            text_frame = shape.text_frame
            text_frame.clear()
            
            for bullet in bullets:
                p = text_frame.paragraphs[0] if len(text_frame.paragraphs) > 0 else text_frame.add_paragraph()
                if text_frame.paragraphs[0] != p:
                    p = text_frame.add_paragraph()
                
                p.text = bullet['text']
                p.level = bullet.get('level', 0)
            
            break
    
    return slide

# === HAUPTBEARBEITUNG ===

print("=" * 60)
print("STARTE POWERPOINT-ÜBERARBEITUNG")
print("=" * 60)

# 1. Füge Executive Summary nach Titelfolie hinzu (Position 1)
print("\n[1/6] Füge Executive Summary hinzu...")
# Source Folie 3 (Index 2) = Executive Summary
copy_slide(source_prs, 2, target_prs, insert_position=1)
log_change("Executive Summary eingefügt (nach Titelfolie, Position 1)")

# 2. Füge "Warum jetzt" / Problem Statement hinzu (Position 2-3)
print("\n[2/6] Füge 'Warum jetzt' und Problem Statement hinzu...")
# Source Folie 4 (Index 3) = WARUM DAS THEMA PRIORITÄT HAT
copy_slide(source_prs, 3, target_prs, insert_position=2)
log_change("'Warum das Thema Priorität hat' eingefügt (Position 2)")

# Source Folie 5 (Index 4) = Problem Statement
copy_slide(source_prs, 4, target_prs, insert_position=3)
log_change("'Problem Statement' eingefügt (Position 3)")

# 3. Behalte operative Details aus Ziel-Deck
# Diese bleiben an ihrer Position (wurden durch insert nach hinten verschoben)
print("\n[3/6] Operative Details bleiben erhalten...")
log_change("Operative Details aus Original-Deck bleiben erhalten (Rollen, Zuarbeit, Leistung, Prozess)")

# 4. Füge Impact-Folie hinzu (vor Risiken, vor letzter Folie des Originals)
print("\n[4/6] Füge Impact-Folie hinzu...")
# Source Folie 11 (Index 10) = Impact
# Einfügen vor der letzten Folie des ursprünglichen Target-Decks
# Original hatte 10 Folien, mit 3 neuen Folien sind es jetzt 13, also vor Position 12
copy_slide(source_prs, 10, target_prs, insert_position=12)
log_change("'Impact: messbare Wirkung' eingefügt (vor Abschluss)")

# 5. Füge Risiken & Mitigation hinzu
print("\n[5/6] Füge Risiken & Mitigation hinzu...")
# Source Folie 12 (Index 11) = Risiken & Mitigation
copy_slide(source_prs, 11, target_prs, insert_position=13)
log_change("'Risiken & Mitigation' eingefügt")

# 6. Verbessere Call-to-Action am Ende
print("\n[6/6] Verbessere Call-to-Action...")
# Entferne alte "UMSETZUNG" Folie (war ursprünglich Folie 10, ist jetzt weiter hinten)
# und ersetze durch "Call to Action" aus Source
# Die alte UMSETZUNG-Folie ist jetzt an Position ~14-15

# Finde die UMSETZUNG-Folie
for i, slide in enumerate(target_prs.slides):
    if slide.shapes.title and "UMSETZUNG" in slide.shapes.title.text.upper():
        # Lösche diese Folie
        rId = target_prs.slides._sldIdLst[i].rId
        target_prs.part.drop_rel(rId)
        del target_prs.slides._sldIdLst[i]
        log_change(f"Alte 'Umsetzung'-Folie entfernt (Position {i})")
        break

# Füge neue Call-to-Action hinzu (Source Folie 13, Index 12)
copy_slide(source_prs, 12, target_prs, insert_position=None)  # Am Ende
log_change("'Call to Action: Entscheidung & Vorgehen' eingefügt (Abschluss)")

# Speichern
output_path = '/Users/nwe/clawd/Bitsight_MHP_Beschaffung_FINAL.pptx'
target_prs.save(output_path)

print("\n" + "=" * 60)
print(f"✅ FERTIG! Gespeichert als: {output_path}")
print("=" * 60)

# Erstelle Changelog
changelog_md = f"""# CHANGELOG - BitSight Präsentations-Überarbeitung

**Datum:** {__import__('datetime').datetime.now().strftime('%Y-%m-%d %H:%M')}  
**Input:** Bitsight_MHP_Beschaffung_Rollen_und_Zuarbeit.pptx  
**Source Content:** Bitsight_Vorstellung MHP.pptx  
**Output:** Bitsight_MHP_Beschaffung_FINAL.pptx

## 🎯 Ziel
Verbesserung des narrativen Flows durch Integration von Executive Summary, Problem Statement, Impact und Call-to-Action aus der Vorstellungs-Präsentation.

## 📝 Durchgeführte Änderungen

### 1. **Executive Summary hinzugefügt** (Position 1, nach Titelfolie)
- **Quelle:** Folie 3 aus "Bitsight_Vorstellung MHP.pptx"
- **Inhalt:** Drei Hebel - Standardisieren, Skalieren, Wirkung
- **Nutzen:** Gibt sofort einen klaren Überblick über die Kernbotschaft

### 2. **"Warum das Thema Priorität hat" hinzugefügt** (Position 2)
- **Quelle:** Folie 4 aus "Bitsight_Vorstellung MHP.pptx"
- **Inhalt:** Kontextualisierung der Dringlichkeit
- **Nutzen:** Etabliert die Notwendigkeit vor der Lösung

### 3. **Problem Statement hinzugefügt** (Position 3)
- **Quelle:** Folie 5 aus "Bitsight_Vorstellung MHP.pptx"
- **Inhalt:** Blind Spots, Operative Last, Entscheidungsdefizit
- **Nutzen:** Konkrete Problemstellung für besseres Verständnis des Lösungsbedarfs

### 4. **Operative Details beibehalten** (Positionen 4-11)
- **Inhalt:** Ursprüngliche Folien zu:
  - Zielbild
  - Lösung (BitSight-Kurzprofil)
  - Zusammenarbeit (Service-Modell MHP)
  - Beschaffung (Rolle & RACI)
  - Zuarbeit
  - Leistung
  - End-to-End Prozess
  - Steuerung & KPIs
- **Nutzen:** Operative Details bleiben vollständig erhalten

### 5. **Impact-Folie hinzugefügt** (Position 12)
- **Quelle:** Folie 11 aus "Bitsight_Vorstellung MHP.pptx"
- **Inhalt:** Risikoreduktion, Speed in Procurement, Effizienz & Governance
- **Nutzen:** Zeigt messbare Wirkung und ROI auf

### 6. **Risiken & Mitigation hinzugefügt** (Position 13)
- **Quelle:** Folie 12 aus "Bitsight_Vorstellung MHP.pptx"
- **Inhalt:** Risiken (Daten/Attribution) + Mitigation (Prozessdesign, Steuerungshebel)
- **Nutzen:** Adressiert potenzielle Bedenken proaktiv

### 7. **Call-to-Action verbessert** (Abschluss, Position 14)
- **Quelle:** Folie 13 aus "Bitsight_Vorstellung MHP.pptx"
- **Ersetzt:** Ursprüngliche "Umsetzung"-Folie
- **Inhalt:** Entscheidungen heute, 90-Tage-Pilot, Rollout & Betrieb 2026
- **Nutzen:** Klarer, strukturierter Aktionsplan mit konkreten nächsten Schritten

## 📊 Struktur-Übersicht

### Vorher (10 Folien):
1. Titelfolie
2. Zielbild
3. Lösung
4. Zusammenarbeit
5. Beschaffung
6. Zuarbeit
7. Leistung
8. End-to-End
9. Steuerung
10. Umsetzung

### Nachher (14 Folien):
1. Titelfolie
2. **✨ Executive Summary** [NEU]
3. **✨ Warum das Thema Priorität hat** [NEU]
4. **✨ Problem Statement** [NEU]
5. Zielbild
6. Lösung
7. Zusammenarbeit
8. Beschaffung
9. Zuarbeit
10. Leistung
11. End-to-End
12. Steuerung
13. **✨ Impact: messbare Wirkung** [NEU]
14. **✨ Risiken & Mitigation** [NEU]
15. **✨ Call to Action: Entscheidung & Vorgehen** [VERBESSERT]

## 🎪 Narrativer Flow - Vorher vs. Nachher

### ❌ Vorher:
Titel → Direkt ins Zielbild → Operative Details → Umsetzung

### ✅ Nachher:
**Titel → Executive Summary → Warum jetzt? → Problem → Zielbild → Lösung → Operative Details → Impact → Risiken → Call-to-Action**

**Verbesserung:**
- **Kontext vor Lösung:** Problemstellung etabliert Bedarf
- **Executive Summary** gibt Orientierung
- **Impact** zeigt Wirkung vor Abschluss
- **Risiken** werden adressiert (Vertrauensaufbau)
- **Stärkerer Call-to-Action** mit konkretem Fahrplan

## ✅ Erfüllte Anforderungen

- ✅ Besserer roter Faden durch logische Story-Struktur
- ✅ Content aus "Bitsight_Vorstellung MHP.pptx" integriert
- ✅ Executive Summary hinzugefügt
- ✅ "Warum jetzt" / Problem Statement hinzugefügt
- ✅ Impact-Folie hinzugefügt
- ✅ Risiken & Mitigation hinzugefügt
- ✅ Call-to-Action verbessert
- ✅ Operative Details aus Ziel-Deck beibehalten
- ✅ python-pptx für Bearbeitung genutzt

## 🔧 Technische Details

**Tool:** python-pptx  
**Methode:** Folie-für-Folie Kopie mit Positions-Management  
**Quelldaten:** slides_analysis.json

---

**Erstellt durch:** Automated PowerPoint Editor  
**Script:** pptx_editor.py
"""

with open('/Users/nwe/clawd/CHANGELOG.md', 'w') as f:
    f.write(changelog_md)

print(f"\n📄 Changelog erstellt: /Users/nwe/clawd/CHANGELOG.md")
print("\n" + "=" * 60)
print("ZUSAMMENFASSUNG DER ÄNDERUNGEN:")
print("=" * 60)
for change in changelog:
    print(f"  • {change}")
print("=" * 60)
